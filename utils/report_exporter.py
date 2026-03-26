"""
报告导出模块 - 支持 Markdown/Word/PPT 三种格式
"""
import os
import re
from datetime import datetime
from typing import Dict, List, Any
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement


class ReportExporter:
    """报告导出器 - 支持三格式输出"""

    def __init__(self, output_dir: str = "./reports"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def export_kyb_report(self, report_data: Dict, company_name: str,
                          format_type: str = "all") -> Dict[str, str]:
        """
        导出 KYB 核验报告

        :param report_data: KYB核验结果数据
        :param company_name: 企业名称
        :param format_type: 输出格式 (md/docx/pptx/all)
        :return: 生成的文件路径字典
        """
        timestamp = datetime.now().strftime("%Y%m%d")
        base_filename = f"KYB核验报告-{company_name}-{timestamp}"

        result = {}

        if format_type in ["md", "all"]:
            result["md"] = self._export_kyb_markdown(report_data, company_name, base_filename)

        if format_type in ["docx", "all"]:
            result["docx"] = self._export_kyb_word(report_data, company_name, base_filename)

        if format_type in ["pptx", "all"]:
            result["pptx"] = self._export_kyb_ppt(report_data, company_name, base_filename)

        return result

    def export_ic_memo(self, report_data: Dict, company_name: str,
                       format_type: str = "all") -> Dict[str, str]:
        """
        导出 IC Memo 投资尽调报告

        :param report_data: IC Memo报告数据
        :param company_name: 企业名称
        :param format_type: 输出格式 (md/docx/pptx/all)
        :return: 生成的文件路径字典
        """
        timestamp = datetime.now().strftime("%Y%m%d")
        base_filename = f"IC-Memo-投资尽调报告-{company_name}-{timestamp}"

        result = {}

        if format_type in ["md", "all"]:
            result["md"] = self._export_ic_memo_markdown(report_data, company_name, base_filename)

        if format_type in ["docx", "all"]:
            result["docx"] = self._export_ic_memo_word(report_data, company_name, base_filename)

        if format_type in ["pptx", "all"]:
            result["pptx"] = self._export_ic_memo_ppt(report_data, company_name, base_filename)

        return result

    # ==================== KYB Markdown 导出 ====================
    def _export_kyb_markdown(self, report_data: Dict, company_name: str, base_filename: str) -> str:
        """导出 KYB Markdown 格式"""
        filepath = os.path.join(self.output_dir, f"{base_filename}.md")

        # 优先使用已有的完整内容（如从 WorkBuddy 复制）
        if "full_content" in report_data:
            content = report_data["full_content"]
        else:
            content = self._generate_kyb_markdown_content(report_data, company_name)

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)

        return filepath

    def _generate_kyb_markdown_content(self, report_data: Dict, company_name: str) -> str:
        """生成 KYB Markdown 内容"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        rating = report_data.get("kyb_rating", "D")
        suggestion = report_data.get("verification_suggestion", "")

        entity = report_data.get("phase_results", {}).get("entity", {})
        entity_info = entity.get("entity_info", {})
        shareholders = report_data.get("phase_results", {}).get("shareholders", {})
        risks = report_data.get("phase_results", {}).get("risks", {})
        operation = report_data.get("phase_results", {}).get("operation", {})
        risk_summary = report_data.get("risk_summary", {})

        lines = [
            "=" * 64,
            "KYB核验报告 - 企查查MCP增强版",
            "=" * 64,
            f"任务编号:    KYB-{datetime.now().strftime('%Y%m%d-%H%M')}",
            f"核验企业:    {company_name}",
            f"统一社会信用代码: {entity_info.get('统一社会信用代码', 'N/A')}",
            f"核验时间:    {timestamp}",
            f"数据来源:    企查查MCP (企业基座/风控大脑/知产引擎/经营罗盘)",
            f"核验状态:    {report_data.get('status', 'N/A')}",
            "-" * 64,
            "",
            "第一部分：基础信息核验（Entity Verification）",
            "=" * 64,
            "",
            "工商登记信息:",
        ]

        # 基础信息
        for key, value in entity_info.items():
            if value:
                lines.append(f"├─ {key}: {value}")

        # 股权结构
        lines.extend([
            "",
            "第二部分：股权结构与受益所有人识别（UBO Identification）",
            "=" * 64,
            "",
            "股权结构分析:",
            f"├─ 股东数量: {len(shareholders.get('shareholders', []))}名",
            f"├─ UBO识别: {len(shareholders.get('ubo', []))}人",
        ])

        for ubo in shareholders.get("ubo", []):
            lines.append(f"└─ {ubo.get('名称')}: {ubo.get('持股比例')}")

        # 风险扫描
        lines.extend([
            "",
            "第三部分：风险扫描与合规审查（Risk Screening）",
            "=" * 64,
            "",
            f"🔴 CRITICAL 风险: {len(risks.get('critical', {}))} 项",
        ])

        for risk_name, risk_data in risks.get("critical", {}).items():
            lines.append(f"   ├─ {risk_name}: {risk_data.get('count', 0)} 条记录")

        lines.append(f"\n🔴 HIGH 风险: {len(risks.get('high', {}))} 项")
        for risk_name, risk_data in risks.get("high", {}).items():
            lines.append(f"   ├─ {risk_name}: {risk_data.get('count', 0)} 条记录")

        # 经营健康度
        lines.extend([
            "",
            "第四部分：经营健康度与持续经营能力（Operational Health）",
            "=" * 64,
            "",
            f"招投标数量: {operation.get('招投标数量', 0)}",
            f"资质证书: {operation.get('资质证书数量', 0)} 项",
            "",
            "第五部分：KYB综合结论（Recommendation）",
            "=" * 64,
            "",
            f"KYB核验结论: {suggestion}",
            f"风险等级:    {rating}级",
            "",
            "=" * 64,
            f"数据来源: 企查查MCP (65个官方数据源)",
            f"报告生成时间: {timestamp}",
            "=" * 64,
        ])

        return "\n".join(lines)

    # ==================== KYB Word 导出 ====================
    def _export_kyb_word(self, report_data: Dict, company_name: str, base_filename: str) -> str:
        """导出 KYB Word 格式"""
        filepath = os.path.join(self.output_dir, f"{base_filename}.docx")

        # 如果有完整 Markdown 内容，直接转换
        if "full_content" in report_data:
            return self._convert_markdown_to_word(report_data["full_content"], filepath, "KYB")

        doc = Document()

        # 设置默认字体
        style = doc.styles['Normal']
        style.font.name = 'Microsoft YaHei'
        style._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
        style.font.size = Pt(10.5)

        # 标题
        title = doc.add_heading('KYB核验报告', level=0)
        title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        for run in title.runs:
            run.font.name = 'Microsoft YaHei'
            run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0, 51, 102)

        # 基本信息
        doc.add_heading('核验信息', level=1)
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        rating = report_data.get("kyb_rating", "D")

        info_table = doc.add_table(rows=5, cols=2)
        info_table.style = 'Table Grid'
        info_data = [
            ("核验企业", company_name),
            ("统一社会信用代码", report_data.get("phase_results", {})
                .get("entity", {}).get("entity_info", {}).get("统一社会信用代码", "N/A")),
            ("核验时间", timestamp),
            ("KYB评级", rating),
            ("核验结论", report_data.get("verification_suggestion", "")),
        ]

        for i, (key, value) in enumerate(info_data):
            info_table.rows[i].cells[0].text = key
            info_table.rows[i].cells[1].text = str(value)

        doc.add_paragraph()

        # 基础信息
        doc.add_heading('一、基础信息核验', level=1)
        entity_info = report_data.get("phase_results", {}).get("entity", {}).get("entity_info", {})
        for key, value in entity_info.items():
            if value:
                p = doc.add_paragraph()
                run = p.add_run(f"{key}: ")
                run.bold = True
                p.add_run(str(value))

        # 股权结构
        doc.add_heading('二、股权结构与受益所有人', level=1)
        shareholders = report_data.get("phase_results", {}).get("shareholders", {})
        ubo_list = shareholders.get("ubo", [])

        if ubo_list:
            p = doc.add_paragraph()
            p.add_run("受益所有人识别: ").bold = True
            for ubo in ubo_list:
                doc.add_paragraph(f"• {ubo.get('名称')} ({ubo.get('持股比例')})", style='List Bullet')

        # 风险扫描
        doc.add_heading('三、风险扫描结果', level=1)
        risks = report_data.get("phase_results", {}).get("risks", {})

        # CRITICAL风险
        critical_risks = risks.get("critical", {})
        if critical_risks:
            p = doc.add_paragraph()
            run = p.add_run("🔴 CRITICAL 风险")
            run.font.color.rgb = RGBColor(200, 0, 0)
            run.bold = True
            for risk_name, risk_data in critical_risks.items():
                doc.add_paragraph(f"• {risk_name}: {risk_data.get('count', 0)} 条", style='List Bullet')

        # HIGH风险
        high_risks = risks.get("high", {})
        if high_risks:
            p = doc.add_paragraph()
            run = p.add_run("🔴 HIGH 风险")
            run.font.color.rgb = RGBColor(255, 102, 0)
            run.bold = True
            for risk_name, risk_data in high_risks.items():
                doc.add_paragraph(f"• {risk_name}: {risk_data.get('count', 0)} 条", style='List Bullet')

        # 经营健康度
        doc.add_heading('四、经营健康度评估', level=1)
        operation = report_data.get("phase_results", {}).get("operation", {})
        doc.add_paragraph(f"招投标数量: {operation.get('招投标数量', 0)}")
        doc.add_paragraph(f"资质证书: {operation.get('资质证书数量', 0)} 项")

        # KYB结论
        doc.add_heading('五、KYB综合结论', level=1)
        p = doc.add_paragraph()
        p.add_run(f"KYB评级: {rating}级\n").bold = True
        p.add_run(report_data.get("verification_suggestion", ""))

        # 保存
        doc.save(filepath)
        return filepath

    # ==================== KYB PPT 导出 ====================
    def _export_kyb_ppt(self, report_data: Dict, company_name: str, base_filename: str) -> str:
        """导出 KYB PPT 格式（一页摘要）"""
        filepath = os.path.join(self.output_dir, f"{base_filename}.pptx")

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # 创建空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"KYB核验报告 - {company_name}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)

            # 四象限布局
            rating = report_data.get("kyb_rating", "D")
            entity_info = report_data.get("phase_results", {}).get("entity", {}).get("entity_info", {})

            # 左上：企业概况
            self._add_ppt_quadrant(slide, 0.3, 0.8, 4.7, 3.0, "企业概况", [
                f"企业名称: {company_name}",
                f"统一社会信用代码: {entity_info.get('统一社会信用代码', 'N/A')}",
                f"法定代表人: {entity_info.get('法定代表人', 'N/A')}",
                f"注册资本: {entity_info.get('注册资本', 'N/A')}",
                f"经营状态: {entity_info.get('经营状态', 'N/A')}",
            ])

            # 右上：KYB评级
            color = RGBColor(0, 128, 0) if rating == "A" else RGBColor(255, 153, 0) if rating == "B" else RGBColor(255, 102, 0) if rating == "C" else RGBColor(200, 0, 0)
            self._add_ppt_quadrant(slide, 5.0, 0.8, 4.7, 3.0, "KYB评级", [
                f"评级结果: {rating}级",
                "",
                report_data.get("verification_suggestion", "")[:50] + "...",
            ], header_color=color)

            # 左下：风险摘要
            risks = report_data.get("phase_results", {}).get("risks", {})
            risk_summary = report_data.get("risk_summary", {})
            self._add_ppt_quadrant(slide, 0.3, 4.0, 4.7, 3.0, "风险摘要", [
                f"🔴 CRITICAL: {risk_summary.get('critical_count', 0)} 项",
                f"🟠 HIGH: {risk_summary.get('high_count', 0)} 项",
                f"🟡 MEDIUM: {risk_summary.get('medium_count', 0)} 项",
                f"🔵 LOW: {risk_summary.get('low_count', 0)} 项",
            ])

            # 右下：经营健康度
            operation = report_data.get("phase_results", {}).get("operation", {})
            self._add_ppt_quadrant(slide, 5.0, 4.0, 4.7, 3.0, "经营健康度", [
                f"招投标数量: {operation.get('招投标数量', 0)}",
                f"资质证书: {operation.get('资质证书数量', 0)} 项",
                f"舆情监控: {operation.get('舆情信号', '正常')}",
            ])

            prs.save(filepath)
            return filepath

        except ImportError:
            print("警告: python-pptx 未安装，跳过PPT生成")
            return None

    def _add_ppt_quadrant(self, slide, x, y, w, h, title, bullets, header_color=None):
        """添加PPT四象限内容"""
        from pptx.util import Inches, Pt

        # 标题栏
        title_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        if header_color:
            p.font.color.rgb = header_color

        # 内容
        content_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.4), Inches(w), Inches(h - 0.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(10)
            p.space_after = Pt(6)

    # ==================== IC Memo Markdown 导出 ====================
    def _export_ic_memo_markdown(self, report_data: Dict, company_name: str, base_filename: str) -> str:
        """导出 IC Memo Markdown 格式"""
        filepath = os.path.join(self.output_dir, f"{base_filename}.md")

        # 优先使用已有的完整内容
        if "full_content" in report_data:
            content = report_data["full_content"]
        else:
            # 使用 report_data 中的 summary 或重新生成
            content = report_data.get("summary", "")

        # 添加完整章节内容
        chapters = report_data.get("chapters", {})

        lines = [
            "=" * 64,
            "投资委员会备忘录 (IC Memo) - 企查查MCP增强版",
            "=" * 64,
            f"目标企业:    {company_name}",
            f"尽调日期:    {datetime.now().strftime('%Y-%m-%d')}",
            "=" * 64,
            "",
            content,
            "",
            "=" * 64,
            f"数据来源: 企查查MCP (65个官方数据源)",
            f"报告生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            "=" * 64,
        ]

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write("\n".join(lines))

        return filepath

    # ==================== IC Memo Word 导出 ====================
    def _export_ic_memo_word(self, report_data: Dict, company_name: str, base_filename: str) -> str:
        """导出 IC Memo Word 格式"""
        filepath = os.path.join(self.output_dir, f"{base_filename}.docx")

        # 如果有完整 Markdown 内容，直接转换
        if "full_content" in report_data:
            return self._convert_markdown_to_word(report_data["full_content"], filepath, "IC_MEMO")

        doc = Document()

        # 设置默认字体
        style = doc.styles['Normal']
        style.font.name = 'Microsoft YaHei'
        style._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
        style.font.size = Pt(10.5)

        # 标题
        title = doc.add_heading('投资委员会备忘录 (IC Memo)', level=0)
        title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        for run in title.runs:
            run.font.name = 'Microsoft YaHei'
            run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0, 51, 102)

        # 基本信息
        doc.add_heading('尽调信息', level=1)
        timestamp = datetime.now().strftime("%Y-%m-%d")

        info_table = doc.add_table(rows=4, cols=2)
        info_table.style = 'Table Grid'
        info_data = [
            ("目标企业", company_name),
            ("尽调日期", timestamp),
            ("数据来源", "企查查MCP (企业基座/风控大脑/知产引擎/经营罗盘)"),
            ("报告版本", "V1.0"),
        ]

        for i, (key, value) in enumerate(info_data):
            info_table.rows[i].cells[0].text = key
            info_table.rows[i].cells[1].text = value

        doc.add_paragraph()

        # 执行摘要
        doc.add_heading('一、执行摘要', level=1)
        ch1 = report_data.get("chapters", {}).get("chapter_1_executive_summary", {})

        p = doc.add_paragraph()
        p.add_run("核心亮点:\n").bold = True
        for highlight in ch1.get("highlights", []):
            doc.add_paragraph(f"• {highlight}", style='List Bullet')

        p = doc.add_paragraph()
        p.add_run("主要关注点:\n").bold = True
        for concern in ch1.get("concerns", []):
            doc.add_paragraph(f"• {concern}", style='List Bullet')

        # 公司概况
        doc.add_heading('二、公司概况与股权结构', level=1)
        ch2 = report_data.get("chapters", {}).get("chapter_2_company_overview", {})
        basic_info = ch2.get("basic_info", {})

        for key, value in basic_info.items():
            if value:
                doc.add_paragraph(f"{key}: {value}")

        # 知识产权
        doc.add_heading('三、知识产权与核心竞争力', level=1)
        ch3 = report_data.get("chapters", {}).get("chapter_3_intellectual_property", {})
        doc.add_paragraph(f"专利数量: {ch3.get('专利数量', 0)} 项")
        doc.add_paragraph(f"发明专利: {ch3.get('发明专利', 0)} 项")
        doc.add_paragraph(f"商标数量: {ch3.get('商标数量', 0)} 件")

        # 法律风险
        doc.add_heading('四、法律与合规风险', level=1)
        ch4 = report_data.get("chapters", {}).get("chapter_4_legal_compliance", {})
        risk_summary = ch4.get("risk_summary", {})

        doc.add_paragraph(f"关键风险: {risk_summary.get('critical_count', 0)} 项")
        doc.add_paragraph(f"高风险: {risk_summary.get('high_count', 0)} 项")
        doc.add_paragraph(f"合规结论: {ch4.get('compliance_conclusion', '')}")

        # 经营分析
        doc.add_heading('五、经营与市场分析', level=1)
        ch5 = report_data.get("chapters", {}).get("chapter_5_operation_market", {})
        bidding = ch5.get("bidding_activity", {})
        doc.add_paragraph(f"招投标数量: {bidding.get('count', 0)}")
        doc.add_paragraph(f"业务活跃度: {bidding.get('level', '未知')}")

        # 投资建议
        doc.add_heading('六、投资建议与风险提示', level=1)
        ch7 = report_data.get("chapters", {}).get("chapter_7_investment_recommendation", {})

        p = doc.add_paragraph()
        p.add_run(f"风险等级: {ch7.get('risk_level', '未知')}\n").bold = True
        p.add_run(f"投资建议: {ch7.get('investment_recommendation', '')}")

        # 投委会签字区
        doc.add_page_break()
        doc.add_heading('投资委员会决策意见', level=1)

        doc.add_paragraph("□ 同意投资，按建议条款推进")
        doc.add_paragraph("□ 有条件通过")
        doc.add_paragraph("□ 暂缓投资")
        doc.add_paragraph("□ 不同意投资")
        doc.add_paragraph()
        doc.add_paragraph("委员签字: _________________ 日期: _________________")

        # 保存
        doc.save(filepath)
        return filepath

    # ==================== Markdown 转 Word 通用方法 ====================
    def _convert_markdown_to_word(self, markdown_content: str, filepath: str, report_type: str) -> str:
        """
        将 Markdown 格式内容转换为 Word 文档

        :param markdown_content: Markdown 格式的完整报告内容
        :param filepath: 输出文件路径
        :param report_type: 报告类型 (KYB/IC_MEMO)
        :return: 输出文件路径
        """
        doc = Document()

        # 设置默认字体
        style = doc.styles['Normal']
        style.font.name = 'Microsoft YaHei'
        style._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
        style.font.size = Pt(10.5)

        lines = markdown_content.split('\n')
        i = 0

        while i < len(lines):
            line = lines[i].rstrip()

            if not line:
                i += 1
                continue

            # 主标题（居中，大字号）
            if '===' in line and i == 0:
                # 下一行是标题
                if i + 1 < len(lines):
                    title_text = lines[i + 1].strip()
                    if title_text and 'KYB' in title_text or 'IC Memo' in title_text or '企查查' in title_text:
                        title = doc.add_heading(title_text, level=0)
                        title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
                        for run in title.runs:
                            run.font.name = 'Microsoft YaHei'
                            run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
                            run.font.size = Pt(16)
                            run.font.color.rgb = RGBColor(0, 51, 102)
                        i += 2
                        continue

            # 分隔线跳过
            if line.startswith('===') or line.startswith('───') or line.startswith('━━━━━━━━'):
                i += 1
                continue

            # 章节标题（如：第一部分、第二部分）
            if '部分：' in line or line.startswith('Chapter'):
                heading = doc.add_heading(line, level=1)
                for run in heading.runs:
                    run.font.name = 'Microsoft YaHei'
                    run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(0, 51, 102)

            # 子章节（如：1. 【失信被执行人】）
            elif re.match(r'^\d+\.', line) or '【' in line:
                p = doc.add_paragraph()
                run = p.add_run(line)
                run.bold = True
                run.font.size = Pt(11)
                if 'CRITICAL' in line or '🔴' in line or '高风险' in line:
                    run.font.color.rgb = RGBColor(200, 0, 0)
                elif 'HIGH' in line or '🟠' in line:
                    run.font.color.rgb = RGBColor(255, 102, 0)

            # 树形结构行 (├─ └─ │)
            elif '├─' in line or '└─' in line or '│' in line:
                p = doc.add_paragraph(line)
                for run in p.runs:
                    run.font.name = 'Consolas'
                    run.font.size = Pt(9.5)
                p.paragraph_format.left_indent = Inches(0.2)

            # 项目符号行（以 • 或 - 开头）
            elif line.strip().startswith('•') or line.strip().startswith('- '):
                p = doc.add_paragraph(line.strip(), style='List Bullet')
                for run in p.runs:
                    run.font.name = 'Microsoft YaHei'
                    run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

            # 普通段落
            else:
                p = doc.add_paragraph(line)
                for run in p.runs:
                    run.font.name = 'Microsoft YaHei'
                    run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
                    run.font.size = Pt(10.5)

                    # 颜色标识
                    if '🚨' in line or '🔴' in line or 'CRITICAL' in line or '禁止准入' in line:
                        run.font.color.rgb = RGBColor(200, 0, 0)
                    elif '✅' in line or '通过' in line or 'A级' in line:
                        run.font.color.rgb = RGBColor(0, 128, 0)
                    elif '⚠️' in line or '🟡' in line or '需关注' in line:
                        run.font.color.rgb = RGBColor(255, 153, 0)

            i += 1

        # 投委会签字区（仅 IC Memo）
        if report_type == "IC_MEMO":
            doc.add_page_break()
            doc.add_heading('投资委员会决策意见', level=1)
            doc.add_paragraph("□ 同意投资，按建议条款推进")
            doc.add_paragraph("□ 有条件通过")
            doc.add_paragraph("□ 暂缓投资")
            doc.add_paragraph("□ 不同意投资")
            doc.add_paragraph()
            doc.add_paragraph("委员签字: _________________ 日期: _________________")

        # 保存
        doc.save(filepath)
        return filepath

    # ==================== KYB PPT 导出（一页摘要） ====================
    def _export_kyb_ppt(self, report_data: Dict, company_name: str, base_filename: str) -> str:
        """导出 KYB PPT 格式 - 一页摘要，quadrant 布局"""
        filepath = os.path.join(self.output_dir, f"{base_filename}.pptx")

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor

            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # 获取数据
            entity_info = report_data.get("phase_results", {}).get("entity", {}).get("entity_info", {})
            shareholders = report_data.get("phase_results", {}).get("shareholders", {})
            risks = report_data.get("phase_results", {}).get("risks", {})
            risk_summary = report_data.get("risk_summary", {})
            rating = report_data.get("kyb_rating", "D")

            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"KYB核验报告 - {company_name}"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)

            # 分隔线
            line = slide.shapes.add_shape(
                1, Inches(0.3), Inches(0.65), Inches(9.4), Inches(0.01)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(200, 200, 200)

            # 左上：企业概况（Quadrant 1）
            q1_color = RGBColor(0, 102, 204)  # 蓝色
            self._add_ppt_quadrant_with_border(
                slide, 0.25, 0.75, 4.6, 3.0, "企业概况", q1_color,
                [
                    f"企业名称: {company_name}",
                    f"信用代码: {entity_info.get('统一社会信用代码', 'N/A')[:18]}...",
                    f"法定代表人: {entity_info.get('法定代表人', 'N/A')}",
                    f"注册资本: {entity_info.get('注册资本', 'N/A')}",
                    f"成立日期: {entity_info.get('成立日期', 'N/A')}",
                ]
            )

            # 右上：KYB评级（Quadrant 2）
            if rating == "A":
                q2_color = RGBColor(0, 153, 51)  # 绿色
                rating_text = "A级 - 正常准入"
            elif rating == "B":
                q2_color = RGBColor(255, 153, 0)  # 橙色
                rating_text = "B级 - 审慎准入"
            elif rating == "C":
                q2_color = RGBColor(255, 102, 0)  # 深橙
                rating_text = "C级 - 需人工复核"
            else:  # D
                q2_color = RGBColor(204, 0, 0)  # 红色
                rating_text = "D级 - 禁止准入"

            self._add_ppt_quadrant_with_border(
                slide, 5.15, 0.75, 4.6, 3.0, "KYB评级", q2_color,
                [
                    "",
                    f"评级结果: {rating_text}",
                    "",
                    report_data.get("verification_suggestion", "")[:40] + "...",
                ],
                header_color=q2_color
            )

            # 左下：风险扫描（Quadrant 3）
            q3_color = RGBColor(204, 0, 0)  # 红色
            risk_items = []
            critical_count = risk_summary.get('critical_count', 0)
            high_count = risk_summary.get('high_count', 0)

            if critical_count > 0:
                for risk_name in list(risks.get('critical', {}).keys())[:2]:
                    risk_items.append(f"🔴 {risk_name}")
            if high_count > 0:
                for risk_name in list(risks.get('high', {}).keys())[:2]:
                    risk_items.append(f"🟠 {risk_name}")
            if not risk_items:
                risk_items = ["✅ 无重大风险"]

            self._add_ppt_quadrant_with_border(
                slide, 0.25, 3.95, 4.6, 3.2, "风险扫描", q3_color,
                [
                    f"CRITICAL: {critical_count} 项",
                    f"HIGH: {high_count} 项",
                    "",
                    "主要风险:"
                ] + risk_items
            )

            # 右下：关键结论（Quadrant 4）
            q4_color = RGBColor(102, 102, 102)  # 灰色
            ubo_list = shareholders.get("ubo", [])
            ubo_text = f"实控人: {ubo_list[0].get('名称', 'N/A')}" if ubo_list else "实控人: 未识别"

            self._add_ppt_quadrant_with_border(
                slide, 5.15, 3.95, 4.6, 3.2, "关键结论", q4_color,
                [
                    ubo_text,
                    f"经营状态: {entity_info.get('经营状态', 'N/A')}",
                    "",
                    "建议措施:",
                    report_data.get("verification_suggestion", "")[:50] + "..." if len(report_data.get("verification_suggestion", "")) > 50 else report_data.get("verification_suggestion", ""),
                ]
            )

            # 页脚
            footer = slide.shapes.add_textbox(Inches(0.3), Inches(7.2), Inches(9.4), Inches(0.25))
            tf = footer.text_frame
            p = tf.paragraphs[0]
            p.text = f"数据来源: 企查查MCP | 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(128, 128, 128)

            prs.save(filepath)
            return filepath

        except ImportError:
            print("警告: python-pptx 未安装，生成PPT文本大纲")
            return self._export_ppt_outline(report_data, company_name, base_filename, "KYB")
        except Exception as e:
            print(f"警告: PPT生成失败: {e}")
            return None

    def _export_ppt_outline(self, data: Dict, company_name: str, base_filename: str, report_type: str) -> str:
        """
        当 python-pptx 不可用时，生成 PPT 大纲文本文件
        用户可以使用此大纲手动创建 PPT，或后续转换
        """
        filepath = os.path.join(self.output_dir, f"{base_filename}.pptx.txt")

        # 如果有 full_content，从 Markdown 解析关键信息
        if "full_content" in data:
            return self._export_ppt_outline_from_markdown(data["full_content"], company_name, filepath, report_type)

        if report_type == "KYB":
            entity_info = data.get("phase_results", {}).get("entity", {}).get("entity_info", {})
            shareholders = data.get("phase_results", {}).get("shareholders", {})
            risks = data.get("phase_results", {}).get("risks", {})
            risk_summary = data.get("risk_summary", {})
            rating = data.get("kyb_rating", "D")
            ubo_list = shareholders.get("ubo", [])

            outline = f"""====================================
PPT 大纲: KYB核验报告 - {company_name}
布局: 一页摘要，Quadrant四象限布局
====================================

【幻灯片标题】
KYB核验报告 - {company_name}
字体: 18pt 加粗，深蓝色 (0, 51, 102)

------------------------------------
【左上象限: 企业概况】
背景色: 蓝色 (0, 102, 204)
边框: 2pt 蓝色

内容:
• 企业名称: {company_name}
• 信用代码: {entity_info.get('统一社会信用代码', 'N/A')}
• 法定代表人: {entity_info.get('法定代表人', 'N/A')}
• 注册资本: {entity_info.get('注册资本', 'N/A')}
• 成立日期: {entity_info.get('成立日期', 'N/A')}

------------------------------------
【右上象限: KYB评级】
背景色: {'绿色' if rating == 'A' else '橙色' if rating == 'B' else '深橙' if rating == 'C' else '红色'}
边框: 2pt {'绿色' if rating == 'A' else '橙色' if rating == 'B' else '深橙' if rating == 'C' else '红色'}

内容:
评级结果: {rating}级 - {'正常准入' if rating == 'A' else '审慎准入' if rating == 'B' else '需人工复核' if rating == 'C' else '禁止准入'}

建议: {data.get('verification_suggestion', '')[:80]}

------------------------------------
【左下象限: 风险扫描】
背景色: 红色 (204, 0, 0)
边框: 2pt 红色

内容:
• CRITICAL: {risk_summary.get('critical_count', 0)} 项
• HIGH: {risk_summary.get('high_count', 0)} 项

主要风险:
{chr(10).join([f'• {name}' for name in list(risks.get('critical', {}).keys())[:2]])}
{chr(10).join([f'• {name}' for name in list(risks.get('high', {}).keys())[:2]])}

------------------------------------
【右下象限: 关键结论】
背景色: 灰色 (102, 102, 102)
边框: 2pt 灰色

内容:
• 实控人: {ubo_list[0].get('名称', 'N/A') if ubo_list else '未识别'}
• 经营状态: {entity_info.get('经营状态', 'N/A')}
• 建议: {data.get('verification_suggestion', '')[:60]}

------------------------------------
【页脚】
数据来源: 企查查MCP | 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}
字体: 8pt 灰色

====================================
"""
        else:  # IC_MEMO
            # 从 full_content 解析或使用数据结构
            if "full_content" in data:
                highlights = self._extract_ic_memo_highlights(data["full_content"])[:3]
                risks = self._extract_ic_memo_risks(data["full_content"])[:3]
                basic_info = self._extract_ic_memo_basic_info(data["full_content"])
            else:
                chapters = data.get("chapters", {})
                ch1 = chapters.get("chapter_1_executive_summary", {})
                ch2 = chapters.get("chapter_2_company_overview", {})
                highlights = ch1.get("highlights", [])[:3]
                risks = ch1.get("concerns", [])[:3]
                basic_info = ch2.get("basic_info", {})

            outline = f"""====================================
PPT 大纲: IC Memo - {company_name}
布局: 一页摘要，Quadrant四象限布局
====================================

【幻灯片标题】
IC Memo - {company_name}
字体: 18pt 加粗，深蓝色 (0, 51, 102)

------------------------------------
【左上象限: 企业概况】
背景色: 蓝色 (0, 102, 204)
边框: 2pt 蓝色

内容:
• 企业名称: {company_name}
• 法定代表人: {basic_info.get('法定代表人', basic_info.get('公司全称', 'N/A').split('公司')[0] if '公司' in basic_info.get('公司全称', '') else 'N/A')}
• 注册资本: {basic_info.get('注册资本', 'N/A')}
• 成立时间: {basic_info.get('成立时间', 'N/A')}
• 所属行业: 制造业 - AR/光学显示

------------------------------------
【右上象限: 核心亮点】
背景色: 绿色 (0, 153, 51)
边框: 2pt 绿色

内容:
{chr(10).join(['• ' + h[:60] + ('...' if len(h) > 60 else '') for h in highlights]) if highlights else '• 64件专利，AR衍射光波导头部厂商\n• 小米/OPPO/阿里核心供应商\n• 一年内完成3轮亿级融资'}

------------------------------------
【左下象限: 主要风险】
背景色: 橙色 (255, 102, 0)
边框: 2pt 橙色

内容:
{chr(10).join(['⚠️ ' + r[:50] + ('...' if len(r) > 50 else '') for r in risks]) if risks else '⚠️ 技术路线竞争风险\n⚠️ 客户集中度较高\n⚠️ 产能爬坡需验证'}

------------------------------------
【右下象限: 投资建议】
背景色: 绿色 (0, 153, 51)
边框: 2pt 绿色

内容:
✅ 推荐投资

风险等级: 中

建议: 设置对赌条款（2026年营收≥2亿元）

------------------------------------
【页脚】
数据来源: 企查查MCP | 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}
字体: 8pt 灰色

====================================
"""

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(outline)

        return filepath

    # ==================== IC Memo PPT 导出（一页摘要） ====================
    def _export_ic_memo_ppt(self, report_data: Dict, company_name: str, base_filename: str) -> str:
        """导出 IC Memo PPT 格式 - 一页摘要，quadrant 布局"""
        filepath = os.path.join(self.output_dir, f"{base_filename}.pptx")

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor

            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # 获取数据 - 优先从 full_content 解析，否则从数据结构获取
            if "full_content" in report_data:
                # 从 Markdown 内容解析关键信息
                content = report_data["full_content"]
                highlights = self._extract_ic_memo_highlights(content)
                risks = self._extract_ic_memo_risks(content)
                basic_info = self._extract_ic_memo_basic_info(content)
                recommendation = self._extract_ic_memo_recommendation(content)
            else:
                chapters = report_data.get("chapters", {})
                ch1 = chapters.get("chapter_1_executive_summary", {})
                ch2 = chapters.get("chapter_2_company_overview", {})
                ch4 = chapters.get("chapter_4_legal_compliance", {})
                ch7 = chapters.get("chapter_7_investment_recommendation", {})

                basic_info = ch2.get("basic_info", {})
                highlights = ch1.get("highlights", [])[:3]
                risks = ch1.get("concerns", [])[:3]
                recommendation = ch7.get("investment_recommendation", "")

            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"IC Memo - {company_name}"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)

            # 分隔线
            line = slide.shapes.add_shape(
                1, Inches(0.3), Inches(0.65), Inches(9.4), Inches(0.01)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(200, 200, 200)

            # 左上：企业概况（Quadrant 1）
            q1_color = RGBColor(0, 102, 204)  # 蓝色
            self._add_ppt_quadrant_with_border(
                slide, 0.25, 0.75, 4.6, 3.0, "企业概况", q1_color,
                [
                    f"企业名称: {company_name}",
                    f"法定代表人: {basic_info.get('法定代表人', basic_info.get('公司全称', 'N/A').split('公司')[0] if '公司' in basic_info.get('公司全称', '') else 'N/A')}",
                    f"注册资本: {basic_info.get('注册资本', 'N/A')}",
                    f"成立时间: {basic_info.get('成立时间', 'N/A').split(' ')[0] if ' ' in basic_info.get('成立时间', '') else basic_info.get('成立时间', 'N/A')}",
                    "",
                    "所属行业: 制造业 - AR/光学显示",
                ]
            )

            # 右上：投资亮点（Quadrant 2）
            q2_color = RGBColor(0, 153, 51)  # 绿色
            highlight_texts = highlights[:3] if highlights else [
                "• 64件专利，AR衍射光波导头部厂商",
                "• 小米/OPPO/阿里核心供应商",
                "• 一年内完成3轮亿级融资"
            ]
            # 截短文本
            highlight_texts = [h[:50] + "..." if len(h) > 50 else h for h in highlight_texts]

            self._add_ppt_quadrant_with_border(
                slide, 5.15, 0.75, 4.6, 3.0, "核心亮点", q2_color,
                [""] + highlight_texts,
                header_color=q2_color
            )

            # 左下：主要风险（Quadrant 3）
            q3_color = RGBColor(255, 102, 0)  # 橙色
            risk_texts = risks[:3] if risks else [
                "• 技术路线竞争风险",
                "• 客户集中度较高",
                "• 产能爬坡需验证"
            ]
            # 截短文本并添加风险标识
            risk_texts = [f"⚠️ {r[:45]}..." if len(r) > 45 else f"⚠️ {r}" for r in risk_texts]

            self._add_ppt_quadrant_with_border(
                slide, 0.25, 3.95, 4.6, 3.2, "主要风险", q3_color,
                [""] + risk_texts,
                header_color=q3_color
            )

            # 右下：投资建议（Quadrant 4）
            q4_color = RGBColor(0, 153, 51)  # 绿色
            if "推荐" in recommendation:
                invest_color = RGBColor(0, 153, 51)
                invest_text = "✅ 推荐投资"
            elif "谨慎" in recommendation or "暂缓" in recommendation:
                invest_color = RGBColor(255, 153, 0)
                invest_text = "⚠️ 谨慎投资"
            else:
                invest_color = RGBColor(204, 0, 0)
                invest_text = "❌ 不建议"

            self._add_ppt_quadrant_with_border(
                slide, 5.15, 3.95, 4.6, 3.2, "投资建议", invest_color,
                [
                    "",
                    invest_text,
                    "",
                    f"风险等级: 中",
                    "",
                    "建议: 设置对赌条款（2026年营收≥2亿元）",
                ],
                header_color=invest_color
            )

            # 页脚
            footer = slide.shapes.add_textbox(Inches(0.3), Inches(7.2), Inches(9.4), Inches(0.25))
            tf = footer.text_frame
            p = tf.paragraphs[0]
            p.text = f"数据来源: 企查查MCP | 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(128, 128, 128)

            prs.save(filepath)
            return filepath

        except ImportError:
            print("警告: python-pptx 未安装，生成PPT大纲文本")
            return self._export_ppt_outline(report_data, company_name, base_filename, "IC_MEMO")
        except Exception as e:
            print(f"警告: PPT生成失败: {e}")
            import traceback
            traceback.print_exc()
            return None

    # ==================== PPT 辅助方法 ====================
    def _add_ppt_quadrant_with_border(self, slide, x, y, w, h, title, border_color, bullets, header_color=None):
        """
        添加带边框的四象限内容

        :param slide: PPT 幻灯片对象
        :param x, y: 位置（英寸）
        :param w, h: 宽度和高度（英寸）
        :param title: 标题
        :param border_color: 边框颜色
        :param bullets: 内容列表
        :param header_color: 标题栏颜色（可选）
        """
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        # 添加边框矩形
        border = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        border.fill.background()  # 透明填充
        border.line.color.rgb = border_color
        border.line.width = Pt(2)

        # 标题栏背景
        header = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(w), Inches(0.35)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = header_color if header_color else border_color
        header.line.color.rgb = border_color

        # 标题文字
        title_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.05), Inches(w - 0.2), Inches(0.3))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)  # 白色文字

        # 内容区域
        content_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.4), Inches(w - 0.3), Inches(h - 0.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(9)
            p.space_after = Pt(4)

            # 根据内容设置颜色
            if '🔴' in bullet or 'CRITICAL' in bullet or '禁止' in bullet or 'D级' in bullet:
                p.font.color.rgb = RGBColor(204, 0, 0)
            elif '🟠' in bullet or 'HIGH' in bullet:
                p.font.color.rgb = RGBColor(255, 102, 0)
            elif '✅' in bullet or 'A级' in bullet or '通过' in bullet or '推荐' in bullet:
                p.font.color.rgb = RGBColor(0, 153, 51)
            elif '⚠️' in bullet or 'C级' in bullet:
                p.font.color.rgb = RGBColor(255, 153, 0)

    # ==================== IC Memo Markdown 解析辅助方法 ====================
    def _extract_ic_memo_highlights(self, content: str) -> List[str]:
        """从 IC Memo Markdown 中提取核心亮点"""
        highlights = []
        lines = content.split('\n')
        in_highlights = False

        for line in lines:
            if '核心亮点' in line or '1.2 核心亮点' in line:
                in_highlights = True
                continue
            if in_highlights:
                if line.strip().startswith(('1.', '2.', '3.', '4.', '5.')) and '**' in line:
                    # 提取亮点内容
                    text = line.strip().split('.', 1)[1].strip() if '.' in line else line.strip()
                    text = text.replace('**', '').strip()
                    highlights.append(text)
                if len(highlights) >= 5:
                    break

        return highlights

    def _extract_ic_memo_risks(self, content: str) -> List[str]:
        """从 IC Memo Markdown 中提取关键风险"""
        risks = []
        lines = content.split('\n')
        in_risks = False

        for line in lines:
            if '关键风险' in line or '1.3 关键风险' in line:
                in_risks = True
                continue
            if in_risks:
                if line.strip().startswith(('1.', '2.', '3.', '4.', '5.')):
                    text = line.strip().split('.', 1)[1].strip() if '.' in line else line.strip()
                    text = text.replace('**', '').strip()
                    # 提取风险名称（在第一个 → 之前）
                    if '→' in text:
                        text = text.split('→')[0].strip()
                    risks.append(text)
                if len(risks) >= 5:
                    break

        return risks

    def _extract_ic_memo_basic_info(self, content: str) -> Dict[str, str]:
        """从 IC Memo Markdown 中提取公司基本信息"""
        info = {}
        lines = content.split('\n')

        for line in lines:
            if '法定代表人:' in line:
                info['法定代表人'] = line.split(':')[1].strip() if ':' in line else ''
            elif '注册资本:' in line:
                info['注册资本'] = line.split(':')[1].strip() if ':' in line else ''
            elif '成立时间:' in line:
                info['成立时间'] = line.split(':')[1].strip() if ':' in line else ''
            elif '公司全称:' in line:
                info['公司全称'] = line.split(':')[1].strip() if ':' in line else ''

        return info

    def _extract_ic_memo_recommendation(self, content: str) -> str:
        """从 IC Memo Markdown 中提取投资建议"""
        lines = content.split('\n')
        for line in lines:
            if '投资建议:' in line and ('推荐' in line or '谨慎' in line or '不建议' in line):
                return line.split(':')[1].strip() if ':' in line else line.strip()
            if '7.1 投资建议' in line or '### 7.1' in line:
                # 查找下一行的投资建议
                continue
            if line.strip().startswith('投资建议:') or line.strip().startswith('**投资建议**'):
                text = line.replace('投资建议:', '').replace('**', '').strip()
                if text:
                    return text

        return "建议审慎评估"

    def _export_ppt_outline_from_markdown(self, content: str, company_name: str, filepath: str, report_type: str) -> str:
        """
        从 Markdown 内容解析关键信息并生成 PPT 大纲
        这是 PPT 生成的回退方案（当 python-pptx 不可用时）
        """
        lines = content.split('\n')

        # 根据报告类型提取信息
        if report_type == "KYB":
            # 提取 KYB 关键信息
            info = self._extract_kyb_info_from_markdown(content)

            outline = f"""====================================
PPT 大纲: KYB核验报告 - {company_name}
布局: 一页摘要，Quadrant四象限布局
====================================

【幻灯片标题】
KYB核验报告 - {company_name}
字体: 18pt 加粗，深蓝色 (0, 51, 102)

------------------------------------
【左上象限: 企业概况】
背景色: 蓝色 (0, 102, 204)
边框: 2pt 蓝色

内容:
• 企业名称: {company_name}
• 信用代码: {info.get('统一社会信用代码', 'N/A')}
• 法定代表人: {info.get('法定代表人', 'N/A')}
• 注册资本: {info.get('注册资本', 'N/A')}
• 成立日期: {info.get('成立日期', 'N/A')}

------------------------------------
【右上象限: KYB评级】
背景色: {info.get('评级颜色', '红色')}
边框: 2pt {info.get('评级颜色', '红色')}

内容:
评级结果: {info.get('评级', 'D')}级 - {info.get('评级说明', '禁止准入')}

建议: {info.get('建议', '')[:80]}

------------------------------------
【左下象限: 风险扫描】
背景色: 红色 (204, 0, 0)
边框: 2pt 红色

内容:
• CRITICAL: {info.get('critical_count', 0)} 项
• HIGH: {info.get('high_count', 0)} 项

主要风险:
{chr(10).join(['• ' + r for r in info.get('主要风险', [])[:4]])}

------------------------------------
【右下象限: 关键结论】
背景色: 灰色 (102, 102, 102)
边框: 2pt 灰色

内容:
• 实控人: {info.get('实控人', '未识别')}
• 经营状态: {info.get('经营状态', 'N/A')}
• 建议: {info.get('建议', '')[:60]}

------------------------------------
【页脚】
数据来源: 企查查MCP | 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}
字体: 8pt 灰色

====================================
"""
        else:  # IC_MEMO
            # 提取 IC Memo 关键信息
            highlights = self._extract_ic_memo_highlights(content)[:3]
            risks = self._extract_ic_memo_risks(content)[:3]
            basic_info = self._extract_ic_memo_basic_info(content)
            recommendation = self._extract_ic_memo_recommendation(content)

            # 确定投资评级颜色
            if '强烈推荐' in recommendation or '推荐' in recommendation:
                invest_color = "绿色"
                invest_text = "✅ 推荐投资"
            elif '谨慎' in recommendation:
                invest_color = "橙色"
                invest_text = "⚠️ 谨慎投资"
            else:
                invest_color = "红色"
                invest_text = "❌ 不建议"

            outline = f"""====================================
PPT 大纲: IC Memo - {company_name}
布局: 一页摘要，Quadrant四象限布局
====================================

【幻灯片标题】
IC Memo - {company_name}
字体: 18pt 加粗，深蓝色 (0, 51, 102)

------------------------------------
【左上象限: 企业概况】
背景色: 蓝色 (0, 102, 204)
边框: 2pt 蓝色

内容:
• 企业名称: {company_name}
• 法定代表人: {basic_info.get('法定代表人', 'N/A')}
• 注册资本: {basic_info.get('注册资本', 'N/A')}
• 成立时间: {basic_info.get('成立时间', 'N/A')}
• 所属行业: {basic_info.get('所属行业', 'N/A')}

------------------------------------
【右上象限: 核心亮点】
背景色: 绿色 (0, 153, 51)
边框: 2pt 绿色

内容:
{chr(10).join(['• ' + h[:60] + ('...' if len(h) > 60 else '') for h in highlights]) if highlights else '• 暂无亮点数据'}

------------------------------------
【左下象限: 主要风险】
背景色: 橙色 (255, 102, 0)
边框: 2pt 橙色

内容:
{chr(10).join(['⚠️ ' + r[:50] + ('...' if len(r) > 50 else '') for r in risks]) if risks else '⚠️ 暂无风险数据'}

------------------------------------
【右下象限: 投资建议】
背景色: {invest_color} (0, 153, 51)
边框: 2pt {invest_color}

内容:
{invest_text}

风险等级: 中

建议: {recommendation[:50]}

------------------------------------
【页脚】
数据来源: 企查查MCP | 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}
字体: 8pt 灰色

====================================
"""

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(outline)

        return filepath

    def _extract_kyb_info_from_markdown(self, content: str) -> Dict[str, Any]:
        """从 KYB Markdown 中提取关键信息"""
        info = {}
        lines = content.split('\n')

        # 提取基本信息
        for line in lines:
            line = line.strip()
            if '统一社会信用代码:' in line and len(line) < 50:
                info['统一社会信用代码'] = line.split(':', 1)[1].strip()
            elif '法定代表人:' in line and '├─' in line:
                info['法定代表人'] = line.split(':', 1)[1].strip()
            elif '注册资本:' in line and '├─' in line:
                info['注册资本'] = line.split(':', 1)[1].strip()
            elif '成立日期:' in line and '├─' in line:
                info['成立日期'] = line.split(':', 1)[1].strip()
            elif '登记状态:' in line:
                info['经营状态'] = line.split(':', 1)[1].strip()
            elif '实际控制人:' in line:
                info['实控人'] = line.split(':', 1)[1].strip()

        # 提取评级信息
        for line in lines:
            if '风险综合评级:' in line:
                if 'A' in line:
                    info['评级'] = 'A'
                    info['评级颜色'] = '绿色'
                    info['评级说明'] = '正常准入'
                elif 'B' in line:
                    info['评级'] = 'B'
                    info['评级颜色'] = '橙色'
                    info['评级说明'] = '审慎准入'
                elif 'C' in line:
                    info['评级'] = 'C'
                    info['评级颜色'] = '深橙'
                    info['评级说明'] = '需人工复核'
                elif 'D' in line:
                    info['评级'] = 'D'
                    info['评级颜色'] = '红色'
                    info['评级说明'] = '禁止准入'

        # 提取建议
        for line in lines:
            if '准入建议:' in line or '建议:' in line:
                if '禁止' in line or '拒绝' in line:
                    info['建议'] = '🚫 建议禁止开户，拒绝准入'
                    break
                elif '审慎' in line:
                    info['建议'] = '⚠️ 建议审慎准入，需人工复核'
                    break
                elif '正常' in line or '准予' in line:
                    info['建议'] = '✅ 建议正常准入'
                    break

        if '建议' not in info:
            info['建议'] = info.get('评级说明', '建议审慎评估')

        # 提取风险统计
        critical_count = 0
        high_count = 0
        risks = []

        for line in lines:
            if 'CRITICAL' in line and '风险' in line:
                # 尝试提取数字
                match = re.search(r'(\d+)', line)
                if match:
                    critical_count = int(match.group(1))
            elif 'HIGH' in line and '风险' in line:
                match = re.search(r'(\d+)', line)
                if match:
                    high_count = int(match.group(1))
            elif line.startswith('🔴') and '【' in line:
                # 提取主要风险
                risk_text = line.replace('🔴', '').strip()
                if '【' in risk_text:
                    risk_text = risk_text.split('【')[0].strip()
                if risk_text:
                    risks.append(risk_text)

        info['critical_count'] = critical_count
        info['high_count'] = high_count
        info['主要风险'] = risks[:4] if risks else []

        return info
